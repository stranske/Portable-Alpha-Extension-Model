from __future__ import annotations

import argparse
from pathlib import Path
import pandas as pd

from pa_core.viz import risk_return, fan, path_dist, corr_heatmap, sharpe_ladder


PLOTS = {
    "risk_return": risk_return.make,
    "fan": fan.make,
    "path_dist": path_dist.make,
    "corr_heatmap": corr_heatmap.make,
    "sharpe_ladder": sharpe_ladder.make,
}


def main(argv: list[str] | None = None) -> None:
    parser = argparse.ArgumentParser(description="Visualise simulation outputs")
    parser.add_argument("--plot", choices=PLOTS.keys(), required=True)
    parser.add_argument("--xlsx", required=True, help="Outputs.xlsx file")
    parser.add_argument("--agent", action="append", help="Agent name for paths")
    parser.add_argument("--png", action="store_true")
    parser.add_argument("--pdf", action="store_true")
    parser.add_argument("--pptx", action="store_true")
    args = parser.parse_args(argv)

    df_summary = pd.read_excel(args.xlsx, sheet_name="Summary")
    parquet_path = Path(args.xlsx).with_suffix(".parquet")
    df_paths = pd.read_parquet(parquet_path) if parquet_path.exists() else None

    if args.plot in {"fan", "path_dist"} and df_paths is None:
        raise FileNotFoundError(parquet_path)

    if args.plot == "corr_heatmap":
        if df_paths is None:
            raise FileNotFoundError(parquet_path)
        paths_map = {"All": df_paths}
        fig = PLOTS[args.plot](paths_map)
    elif args.plot in {"fan", "path_dist"}:
        fig = PLOTS[args.plot](df_paths)
    else:
        fig = PLOTS[args.plot](df_summary)

    base = Path("plots")
    base.mkdir(exist_ok=True)
    stem = base / args.plot
    if args.png:
        fig.write_image(f"{stem}.png")
    if args.pdf:
        fig.write_image(f"{stem}.pdf")
    if args.pptx:
        from pptx import Presentation
        img_path = stem.with_suffix(".png")
        fig.write_image(img_path)
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[5])
        slide.shapes.add_picture(str(img_path), 0, 0)
        pres.save(f"{stem}.pptx")


if __name__ == "__main__":  # pragma: no cover - CLI entry
    main()
