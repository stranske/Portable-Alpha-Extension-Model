"""Tests for issue #1923 — disclose model limitations in docs & board-pack.

Verifies the single source of truth (``pa_core.reporting.disclaimers``) is wired
into both the README and the generated PPTX board pack so the two cannot drift.
"""

import tempfile
from pathlib import Path

import pandas as pd
from pptx import Presentation

from pa_core.reporting.disclaimers import (
    LIMITATIONS_TITLE,
    MODEL_LIMITATIONS,
    limitations_markdown,
)
from pa_core.reporting.export_packet import create_export_packet

_REPO_ROOT = Path(__file__).resolve().parents[1]


def test_limitations_cover_required_caveats():
    """Every caveat named in the audit issue must be present."""
    blob = " ".join(MODEL_LIMITATIONS).lower()
    for phrase in (
        "gross of fees",
        "total excludes base",
        "i.i.d",
        "regimes are ignored",
        "broadcast",
        "not been backtested",
        "risk_metrics",
        "scenario.sleeves",
        "floored at zero",
        "positive carry",
        "single_with_sensitivity",
        "sweep engine",
    ):
        assert phrase in blob, f"missing caveat: {phrase}"


def test_limitations_markdown_is_a_bullet_list():
    md = limitations_markdown()
    assert md.count("- ") == len(MODEL_LIMITATIONS)
    assert MODEL_LIMITATIONS[0] in md


def test_readme_documents_limitations():
    readme = (_REPO_ROOT / "README.md").read_text()
    section_start = readme.index(f"## {LIMITATIONS_TITLE}")
    next_section = readme.index("\n## ", section_start + 1)
    section = readme[section_start:next_section]
    assert limitations_markdown() in section


def test_parameter_guide_documents_advisory_and_unwired_fields():
    guide = (_REPO_ROOT / "docs/guides/PARAMETER_GUIDE.md").read_text()
    for phrase in (
        "`risk_metrics` controls which metrics are reported",
        "`Scenario.sleeves` is validated by the scenario schema",
        "legacy financing costs",
        "floored at zero",
        "`analysis_mode: single_with_sensitivity`",
        "parameter sweep engine supports `returns`, `capital`,",
    ):
        assert phrase in guide


def test_board_pack_includes_limitations_slide():
    """The generated PPTX must carry a limitations slide with every caveat."""
    summary_df = pd.DataFrame(
        {
            "terminal_AnnReturn": [0.05],
            "monthly_AnnVol": [0.12],
            "terminal_ShortfallProb": [0.1],
        }
    )
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path, _ = create_export_packet(
            figs=[],
            summary_df=summary_df,
            raw_returns_dict={"Summary": summary_df},
            inputs_dict={"N_SIMULATIONS": 10, "N_MONTHS": 12},
            base_filename=str(Path(tmpdir) / "packet"),
        )
        prs = Presentation(pptx_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        blob = "\n".join(texts)
        assert LIMITATIONS_TITLE in blob
        for item in MODEL_LIMITATIONS:
            assert item in blob, f"slide missing caveat: {item}"
