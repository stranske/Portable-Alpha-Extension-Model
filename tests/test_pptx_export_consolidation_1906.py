"""Regression tests for issue #1906: consolidate the two PPTX exporters.

`pa_core/viz/pptx_export.py` and `pa_core/reporting/export_packet.py` used to be
near-duplicate PPTX implementations with divergent kaleido error handling (one
silently swallowed render failures, the other raised). Both must now share a
single render/slide code path with consistent, actionable error handling so a
missing renderer can never silently drop chart slides from a board pack.
"""

from __future__ import annotations

from pathlib import Path

import plotly.graph_objects as go
import pytest
from pptx import Presentation

from pa_core.reporting import export_packet
from pa_core.viz import pptx_export


class _BoomFig:
    """Stand-in figure whose static render always fails (no kaleido)."""

    layout = None

    def to_image(self, *args: object, **kwargs: object) -> bytes:
        raise ValueError("Kaleido/Chromium not installed")


class _BadFig:
    """Stand-in figure whose render fails for a non-renderer data reason."""

    layout = None

    def to_image(self, *args: object, **kwargs: object) -> bytes:
        raise ValueError("bad figure data")


def _without_ci_placeholder(monkeypatch: pytest.MonkeyPatch) -> None:
    # Force the real-render branch (pytest normally short-circuits to a
    # placeholder image to avoid spawning kaleido).
    monkeypatch.delenv("CI", raising=False)
    monkeypatch.delenv("PYTEST_CURRENT_TEST", raising=False)
    monkeypatch.delenv("PA_PPTX_PLACEHOLDER", raising=False)


def test_export_packet_delegates_to_shared_helper() -> None:
    # Guard against re-divergence: the committee packet must reuse the viz
    # exporter rather than reimplementing chart rendering.
    assert export_packet.add_chart_slide is pptx_export.add_chart_slide


def test_render_chart_png_uses_placeholder_under_pytest() -> None:
    # Under pytest PYTEST_CURRENT_TEST is set, so no kaleido subprocess runs.
    assert pptx_export.render_chart_png(_BoomFig()) == pptx_export._ONE_PX_PNG


def test_render_chart_png_does_not_mask_renderer_failure_in_generic_ci(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("CI", "1")
    monkeypatch.delenv("PYTEST_CURRENT_TEST", raising=False)
    monkeypatch.delenv("PA_PPTX_PLACEHOLDER", raising=False)

    with pytest.raises(RuntimeError, match="static image renderer"):
        pptx_export.render_chart_png(_BoomFig())


def test_render_chart_png_preserves_non_renderer_errors(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    _without_ci_placeholder(monkeypatch)
    with pytest.raises(ValueError, match="bad figure data"):
        pptx_export.render_chart_png(_BadFig())


def test_save_raises_actionable_error_without_renderer(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    # `save` no longer silently swallows render failures.
    _without_ci_placeholder(monkeypatch)
    with pytest.raises(RuntimeError, match="static image renderer"):
        pptx_export.save([_BoomFig()], tmp_path / "out.pptx")


def test_packet_chart_slide_raises_same_error_without_renderer(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    # The packet path raises the identical actionable error (shared code path).
    _without_ci_placeholder(monkeypatch)
    prs = Presentation()
    with pytest.raises(RuntimeError, match="static image renderer"):
        export_packet._add_chart_slide(prs, _BoomFig())


def test_save_writes_slides_with_alt_text(tmp_path: Path) -> None:
    fig = go.Figure()
    out = tmp_path / "out.pptx"
    pptx_export.save([fig], out, alt_texts=["my chart"])
    assert out.exists()
    pres = Presentation(out)
    # Layout 5 carries a title placeholder, so scan all shapes for the picture.
    descrs = [
        node.get("descr")
        for shape in pres.slides[0].shapes
        for node in shape._element.xpath("./p:nvPicPr/p:cNvPr")
    ]
    assert "my chart" in descrs
