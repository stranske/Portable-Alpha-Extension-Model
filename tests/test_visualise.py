import pandas as pd
from pathlib import Path
from pptx import Presentation

from scripts.visualise import main


def test_visualise_alt_text(tmp_path, monkeypatch):
    summary = pd.DataFrame({
        "AnnReturn": [0.05],
        "AnnVol": [0.02],
        "TrackingErr": [0.01],
        "Agent": ["Base"],
        "ShortfallProb": [0.02],
    })
    xlsx = tmp_path / "out.xlsx"
    summary.to_excel(xlsx, sheet_name="Summary", index=False)
    paths = pd.DataFrame({"Base": [0.01, 0.02]})
    paths.to_parquet(xlsx.with_suffix(".parquet"))
    monkeypatch.chdir(tmp_path)
    main([
        "--plot",
        "risk_return",
        "--xlsx",
        str(xlsx),
        "--html",
        "--pptx",
        "--alt-text",
        "Test chart",
    ])
    html = Path("plots/risk_return.html").read_text()
    assert "aria-label=\"Test chart\"" in html
    pres = Presentation("plots/risk_return.pptx")
    shapes = pres.slides[0].shapes
    elements = shapes[0]._element.xpath('./p:nvPicPr/p:cNvPr')
    if elements:
        assert elements[0].get("descr") == "Test chart"
