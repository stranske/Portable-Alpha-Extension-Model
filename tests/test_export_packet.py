"""Tests for export packet functionality."""

import tempfile
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
import pytest

from pa_core.reporting.export_packet import create_export_packet

# Module-level cached empty DataFrame to avoid repeated creation during tests
# This provides better performance and follows the same pattern used in pa_core/sweep.py
_EMPTY_DATAFRAME = pd.DataFrame()


def create_test_data():
    """Create test data for export packet tests."""
    summary_df = pd.DataFrame(
        {
            "AnnReturn": [0.05, 0.07, 0.06],
            "AnnVol": [0.12, 0.15, 0.13],
            "ShortfallProb": [0.1, 0.15, 0.12],
            "VaR": [0.08, 0.12, 0.10],
        }
    )

    raw_returns_dict = {"Summary": summary_df}

    inputs_dict = {"N_SIMULATIONS": 1000, "N_MONTHS": 12, "analysis_mode": "returns"}

    # Create a simple test figure
    fig = go.Figure()
    fig.add_scatter(
        x=summary_df["AnnVol"],
        y=summary_df["AnnReturn"],
        mode="markers",
        name="Risk-Return",
    )
    fig.update_layout(
        title="Risk-Return Analysis", xaxis_title="Volatility", yaxis_title="Return"
    )

    return summary_df, raw_returns_dict, inputs_dict, fig


def test_export_packet_creates_files():
    """Test that export packet creates both PPTX and Excel files."""
    summary_df, raw_returns_dict, inputs_dict, fig = create_test_data()

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            pptx_path, excel_path = create_export_packet(
                figs=[fig],
                summary_df=summary_df,
                raw_returns_dict=raw_returns_dict,
                inputs_dict=inputs_dict,
                base_filename=str(Path(tmpdir) / "test_packet"),
            )

            # Check that files were created
            assert Path(pptx_path).exists(), f"PPTX file not created: {pptx_path}"
            assert Path(excel_path).exists(), f"Excel file not created: {excel_path}"

            # Check file sizes (should be reasonable)
            pptx_size = Path(pptx_path).stat().st_size
            excel_size = Path(excel_path).stat().st_size

            assert pptx_size > 1000, f"PPTX file too small: {pptx_size} bytes"
            assert excel_size > 1000, f"Excel file too small: {excel_size} bytes"

        except RuntimeError as e:
            if "Chrome" in str(e) or "Chromium" in str(e) or "Kaleido" in str(e):
                pytest.skip("Chrome/Chromium not available in test environment")
            else:
                raise


def test_export_packet_with_alt_text():
    """Test that export packet works with alt text."""
    summary_df, raw_returns_dict, inputs_dict, fig = create_test_data()

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            pptx_path, excel_path = create_export_packet(
                figs=[fig],
                summary_df=summary_df,
                raw_returns_dict=raw_returns_dict,
                inputs_dict=inputs_dict,
                base_filename=str(Path(tmpdir) / "test_packet_alt"),
                alt_texts=["Risk-return scatter plot for investment analysis"],
            )

            # Files should still be created successfully
            assert Path(pptx_path).exists()
            assert Path(excel_path).exists()

        except RuntimeError as e:
            if "Chrome" in str(e) or "Chromium" in str(e) or "Kaleido" in str(e):
                pytest.skip("Chrome/Chromium not available in test environment")
            else:
                raise


def test_export_packet_handles_empty_data():
    """Test that export packet handles edge cases gracefully."""
    # Use module-level cached empty DataFrame instead of creating new instances
    # This follows the recommended pattern and avoids function attribute caching
    empty_summary = _EMPTY_DATAFRAME

    inputs_dict: dict[str, object] = {"test": "value"}
    raw_returns_dict = {"Empty": empty_summary}

    fig = go.Figure()
    fig.update_layout(title="Empty Chart")

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            pptx_path, excel_path = create_export_packet(
                figs=[fig],
                summary_df=empty_summary,
                raw_returns_dict=raw_returns_dict,
                inputs_dict=inputs_dict,
                base_filename=str(Path(tmpdir) / "test_empty"),
            )

            # Should still create files even with empty data
            assert Path(pptx_path).exists()
            assert Path(excel_path).exists()

        except RuntimeError as e:
            if "Chrome" in str(e) or "Chromium" in str(e) or "Kaleido" in str(e):
                pytest.skip("Chrome/Chromium not available in test environment")
            else:
                raise


def test_export_packet_diff_appendix():
    summary_df, raw_returns_dict, inputs_dict, fig = create_test_data()

    prev_summary = summary_df.copy()
    prev_summary["AnnReturn"] = prev_summary["AnnReturn"] + 0.01
    prev_manifest = {"config": {"N_SIMULATIONS": 500}}
    curr_manifest = {"config": {"N_SIMULATIONS": 1000}}

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            pptx_path, excel_path = create_export_packet(
                figs=[fig],
                summary_df=summary_df,
                raw_returns_dict=raw_returns_dict,
                inputs_dict=inputs_dict,
                base_filename=str(Path(tmpdir) / "test_packet_diff"),
                manifest=curr_manifest,
                prev_summary_df=prev_summary,
                prev_manifest=prev_manifest,
            )

            assert Path(pptx_path).exists()
            assert Path(excel_path).exists()

            cfg_diff = pd.read_excel(excel_path, sheet_name="ConfigDiff")
            met_diff = pd.read_excel(excel_path, sheet_name="MetricDiff")
            assert "Parameter" in cfg_diff.columns
            assert "Metric" in met_diff.columns

            from pptx import Presentation

            prs = Presentation(pptx_path)
            texts = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and getattr(shape.text_frame, "text", "")
            ]
            assert any("Config Changes" in t for t in texts)
            assert any("Metric Changes" in t for t in texts)
        except RuntimeError as e:
            if "Chrome" in str(e) or "Chromium" in str(e) or "Kaleido" in str(e):
                pytest.skip("Chrome/Chromium not available in test environment")
            else:
                raise


def test_empty_dataframe_caching():
    """Test that the module-level cached empty DataFrame works as expected."""
    # The cached DataFrame should be the same object when accessed multiple times
    empty1 = _EMPTY_DATAFRAME
    empty2 = _EMPTY_DATAFRAME

    # Same object in memory (identity check)
    assert empty1 is empty2, "Cached DataFrame should be the same object"

    # Proper DataFrame properties
    assert isinstance(empty1, pd.DataFrame), "Should be a DataFrame"
    assert empty1.empty, "Should be empty"
    assert len(empty1) == 0, "Should have no rows"
    assert len(empty1.columns) == 0, "Should have no columns"


def test_no_function_attribute_caching_antipattern():
    """Test that we don't use the function attribute caching anti-pattern.

    This test verifies that we've eliminated the hard-to-maintain pattern:
    if not hasattr(function_name, '_cache'):
        function_name._cache = pd.DataFrame()
    """
    # The function should not have any cache attributes attached to it
    assert not hasattr(
        test_export_packet_handles_empty_data, "_empty_summary_cache"
    ), "Function should not have cache attributes attached"
    assert not hasattr(
        test_export_packet_handles_empty_data, "_cache"
    ), "Function should not have generic cache attributes"

    # Instead, we should use the module-level pattern
    assert "_EMPTY_DATAFRAME" in globals(), "Should use module-level cached variable"
    assert isinstance(
        _EMPTY_DATAFRAME, pd.DataFrame
    ), "Module-level cache should be a DataFrame"


if __name__ == "__main__":
    # Run tests directly if executed as script
    try:
        test_export_packet_creates_files()
        test_export_packet_with_alt_text()
        test_export_packet_handles_empty_data()
        print("✅ All tests passed!")
    except Exception as e:
        print(f"❌ Test failed: {e}")
        raise
