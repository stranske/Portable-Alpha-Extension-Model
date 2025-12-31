import numpy as np
import pandas as pd
import plotly.graph_objects as go

from pa_core.viz import (
    animation,
    beta_heatmap,
    beta_scatter,
    beta_te_scatter,
    bookmark,
    boxplot,
    breach_calendar,
    capital_treemap,
    category_pie,
    corr_heatmap,
    corr_network,
    crossfilter,
    dashboard_templates,
    data_quality,
    data_table,
    delta_heatmap,
    export_bundle,
    exposure_timeline,
    factor_bar,
    factor_matrix,
    factor_timeline,
    fan,
    funnel,
    gauge,
    geo_exposure,
    grid_heatmap,
    grid_panel,
    horizon_slicer,
    hover_sync,
    html_export,
    inset,
    milestone_timeline,
    moments_panel,
    mosaic,
    multi_fan,
    overlay,
    overlay_weighted,
    panel,
    parallel_coords,
    path_dist,
    pdf_export,
    pdf_report,
    pptx_export,
    quantile_band,
    quantile_fan,
    radar,
    rank_table,
    risk_return,
    risk_return_bubble,
    rolling_corr_heatmap,
    rolling_panel,
    rolling_var,
    scatter_matrix,
    scenario_play,
    scenario_viewer,
    seasonality_heatmap,
    sharpe_ladder,
    spark_matrix,
    sunburst,
    surface,
    surface_animation,
    surface_slice,
    te_cvar_scatter,
    triple_scatter,
    violin,
    waterfall,
    weighted_stack,
)


def test_risk_return_make():
    df = pd.DataFrame(
        {
            "AnnReturn": [0.05],
            "AnnVol": [0.02],
            "TrackingErr": [0.01],
            "Agent": ["Base"],
            "ShortfallProb": [0.02],
        }
    )
    fig = risk_return.make(df)
    assert isinstance(fig, go.Figure)
    fig.to_json()


def test_risk_return_axis_labels_match_data():
    df = pd.DataFrame(
        {
            "AnnReturn": [0.06, 0.03],
            "AnnVol": [0.04, 0.02],
            "Agent": ["A", "B"],
            "ShortfallProb": [0.01, 0.02],
        }
    )
    fig = risk_return.make(df)
    assert list(fig.data[0].x) == df["AnnVol"].tolist()
    assert list(fig.data[0].y) == df["AnnReturn"].tolist()
    assert fig.layout.xaxis.title.text == "Annualized Volatility"
    assert fig.layout.yaxis.title.text == "Annualized Return"


def test_fan_and_dist():
    arr = np.random.normal(size=(10, 5))
    fig1 = fan.make(arr)
    fig2 = path_dist.make(arr)
    assert isinstance(fig1, go.Figure)
    assert isinstance(fig2, go.Figure)
    fig1.to_json()
    fig2.to_json()


def test_sharpe_ladder():
    df = pd.DataFrame(
        {
            "AnnReturn": [0.05, 0.03],
            "AnnVol": [0.02, 0.04],
            "Agent": ["A", "B"],
        }
    )
    fig = sharpe_ladder.make(df)
    assert isinstance(fig, go.Figure)
    fig.to_json()


def test_rolling_panel_and_surface_and_pptx(tmp_path):
    arr = np.random.normal(size=(20, 24))
    panel_fig = rolling_panel.make(arr)
    assert isinstance(panel_fig, go.Figure)
    panel_fig.to_json()

    grid = pd.DataFrame(
        {
            "AE_leverage": [1, 1, 2, 2],
            "ExtPA_frac": [0.2, 0.4, 0.2, 0.4],
            "Sharpe": [0.4, 0.5, 0.55, 0.6],
        }
    )
    surf_fig = surface.make(grid)
    assert isinstance(surf_fig, go.Figure)
    surf_fig.to_json()

    out = tmp_path / "out.pptx"
    pptx_export.save([panel_fig, surf_fig], out, alt_texts=["panel", "surface"])
    assert out.exists()
    from pptx import Presentation

    pres = Presentation(out)
    shapes = pres.slides[0].shapes
    assert len(shapes) > 0, "No shapes found on the first slide"
    elements = shapes[0]._element.xpath("./p:nvPicPr/p:cNvPr")
    if elements:
        el = elements[0]
        assert el.get("descr") == "panel"


def test_category_pie():
    fig = category_pie.make(
        {
            "BaseAgent": 500,
            "ExternalPAAgent": 300,
            "InternalPAAgent": 200,
        }
    )
    assert isinstance(fig, go.Figure)
    fig.to_json()


def test_html_and_corr(tmp_path):
    fig = risk_return.make(
        pd.DataFrame(
            {
                "AnnReturn": [0.05],
                "AnnVol": [0.02],
                "TrackingErr": [0.01],
                "Agent": ["Base"],
                "ShortfallProb": [0.02],
            }
        )
    )
    out_html = tmp_path / "fig.html"
    html_export.save(fig, out_html)
    assert out_html.exists()

    arr = np.random.normal(size=(5, 3))
    fig2 = corr_heatmap.make({"A": arr, "B": arr})
    assert isinstance(fig2, go.Figure)
    fig2.to_json()


def test_animation():
    arr = np.random.normal(size=(10, 6))
    fig = animation.make(arr)
    assert isinstance(fig, go.Figure)
    fig.to_json()


def test_overlay_and_waterfall_and_bundle(tmp_path):
    arr = np.random.normal(size=(5, 4))
    over_fig = overlay.make({"A": arr, "B": arr})
    assert isinstance(over_fig, go.Figure)
    over_fig.to_json()

    wf_fig = waterfall.make({"A": 0.1, "B": -0.05})
    assert isinstance(wf_fig, go.Figure)
    wf_fig.to_json()

    export_bundle.save([over_fig, wf_fig], tmp_path / "bundle", alt_texts=["overlay", "waterfall"])
    html = (tmp_path / "bundle_1.html").read_text()
    assert 'aria-label="overlay"' in html


def test_data_table_and_scenario_viewer_and_heatmap():
    df = pd.DataFrame({"A": [1, 2], "B": [3, 4]})
    table_obj = data_table.make(df)
    assert table_obj is not None
    arr = np.random.normal(size=(3, 6))
    fig = scenario_viewer.make({"A": arr})
    assert isinstance(fig, go.Figure)
    fig.to_json()

    grid = pd.DataFrame(
        {
            "AE_leverage": [1, 2],
            "ExtPA_frac": [0.2, 0.4],
            "Sharpe": [0.5, 0.6],
        }
    )
    heatmap_fig = grid_heatmap.make(grid)
    assert isinstance(heatmap_fig, go.Figure)
    heatmap_fig.to_json()


def test_panel_make():
    df = pd.DataFrame(
        {
            "AnnReturn": [0.05, 0.04],
            "AnnVol": [0.02, 0.03],
            "TrackingErr": [0.01, 0.015],
            "Agent": ["A", "B"],
            "ShortfallProb": [0.02, 0.03],
        }
    )
    fig = panel.make(df)
    assert isinstance(fig, go.Figure)
    fig.to_json()


def test_violin_make():
    arr = np.random.normal(size=(10, 5))
    fig = violin.make(arr, by_month=True)
    assert isinstance(fig, go.Figure)
    fig.to_json()


def test_new_viz_helpers():
    arr = np.random.normal(size=(10, 12))
    heat_fig = rolling_corr_heatmap.make(arr, window=3)
    assert isinstance(heat_fig, go.Figure)
    heat_fig.to_json()

    df_cap = pd.DataFrame(
        {
            "A": [100, 120, 110],
            "B": [50, 60, 55],
        },
        index=[0, 1, 2],
    )
    exp_fig = exposure_timeline.make(df_cap)
    assert isinstance(exp_fig, go.Figure)
    exp_fig.to_json()

    summary = pd.DataFrame({"TrackingErr": [0.02]})
    gauge_fig = gauge.make(summary)
    assert isinstance(gauge_fig, go.Figure)
    gauge_fig.to_json()

    metrics = pd.DataFrame(
        {
            "TE": [0.02, 0.03],
            "ER": [0.05, 0.06],
        },
        index=["Scenario1", "Scenario2"],
    )
    radar_fig = radar.make(metrics)
    assert isinstance(radar_fig, go.Figure)
    radar_fig.to_json()


def test_extra_viz_helpers():
    df = pd.DataFrame(
        {
            "AnnReturn": [0.05, 0.04],
            "AnnVol": [0.02, 0.03],
            "TrackingErr": [0.01, 0.015],
            "Capital": [100, 200],
            "ShortfallProb": [0.02, 0.03],
        }
    )
    bubble_fig = risk_return_bubble.make(df)
    assert isinstance(bubble_fig, go.Figure)
    bubble_fig.to_json()

    arr = np.random.normal(size=(10, 12))
    var_fig = rolling_var.make(arr, window=3)
    assert isinstance(var_fig, go.Figure)
    var_fig.to_json()

    summary = pd.DataFrame(
        {
            "Month": [0, 1, 2],
            "TrackingErr": [0.02, 0.04, 0.01],
            "ShortfallProb": [0.01, 0.15, 0.02],
        }
    )
    breach_fig = breach_calendar.make(summary)
    assert isinstance(breach_fig, go.Figure)
    breach_fig.to_json()

    sm_fig = scatter_matrix.make(df[["AnnReturn", "AnnVol", "TrackingErr"]])
    assert isinstance(sm_fig, go.Figure)
    sm_fig.to_json()

    moments_fig = moments_panel.make(arr)
    assert isinstance(moments_fig, go.Figure)
    moments_fig.to_json()

    pc_fig = parallel_coords.make(df[["AnnReturn", "AnnVol", "TrackingErr"]])
    assert isinstance(pc_fig, go.Figure)
    pc_fig.to_json()


def test_newly_added_viz_helpers():
    cap_fig = capital_treemap.make({"A": 100, "B": 200})
    assert isinstance(cap_fig, go.Figure)
    cap_fig.to_json()

    arr = np.random.normal(size=(5, 4))
    net_fig = corr_network.make({"A": arr, "B": arr}, threshold=0.1)
    assert isinstance(net_fig, go.Figure)
    net_fig.to_json()

    beta_df = pd.DataFrame({"A": [1, 2], "B": [0.5, 0.6]}, index=[0, 1])
    beta_fig = beta_heatmap.make(beta_df)
    assert isinstance(beta_fig, go.Figure)
    beta_fig.to_json()


def test_additional_new_viz_helpers():
    df = pd.DataFrame(
        {
            "TrackingErr": [0.02, 0.03],
            "Beta": [1.1, 0.9],
            "Capital": [100, 200],
            "ShortfallProb": [0.01, 0.2],
            "Agent": ["A", "B"],
        }
    )
    bs_fig = beta_scatter.make(df)
    assert isinstance(bs_fig, go.Figure)
    bs_fig.to_json()

    arr = np.random.normal(size=(4, 6))
    ow_fig = overlay_weighted.make({"A": (arr, 1.0)})
    assert isinstance(ow_fig, go.Figure)
    ow_fig.to_json()

    exp_df = pd.DataFrame({"Val": [0.1, 0.2], "Mom": [0.3, -0.1]}, index=["A", "B"])
    fb_fig = factor_bar.make(exp_df)
    assert isinstance(fb_fig, go.Figure)
    fb_fig.to_json()

    fm_fig = factor_matrix.make(exp_df.T)
    assert isinstance(fm_fig, go.Figure)
    fm_fig.to_json()

    mf_fig = multi_fan.make(arr, horizons=[3, 5])
    assert isinstance(mf_fig, go.Figure)
    mf_fig.to_json()


def test_te_cvar_scatter_and_quantile_fan():
    df = pd.DataFrame(
        {
            "TrackingErr": [0.02, 0.03],
            "CVaR": [0.05, 0.04],
            "Agent": ["A", "B"],
        }
    )
    fig1 = te_cvar_scatter.make(df)
    assert isinstance(fig1, go.Figure)
    fig1.to_json()

    arr = np.random.normal(size=(5, 6))
    fig2 = quantile_fan.make(arr, quantiles=(0.1, 0.9))
    assert isinstance(fig2, go.Figure)
    fig2.to_json()


def test_horizon_inset_and_quality(tmp_path):
    arr = np.random.normal(size=(5, 10))
    fig = horizon_slicer.make(arr)
    assert isinstance(fig, go.Figure)
    fig.to_json()

    inset_fig = inset.make(fig, (0, 3, 0.9, 1.1))
    assert isinstance(inset_fig, go.Figure)
    inset_fig.to_json()

    err_df = pd.DataFrame([[1, 0], [0, 2]], index=["2024-01", "2024-02"], columns=["A", "B"])
    qual_fig = data_quality.make(err_df)
    assert isinstance(qual_fig, go.Figure)
    qual_fig.to_json()

    bm = bookmark.save(fig)
    fig2 = bookmark.load(bm)
    assert isinstance(fig2, go.Figure)

    out = tmp_path / "out.pdf"
    pdf_export.save(fig2, out)
    assert out.exists()


def test_sunburst_make():
    df = pd.DataFrame(
        {
            "Agent": ["A", "A", "B"],
            "Sub": ["S1", "S2", "S1"],
            "Return": [0.05, 0.02, -0.01],
        }
    )
    fig = sunburst.make(df)
    assert isinstance(fig, go.Figure)
    fig.to_json()


def test_additional_visualisations(tmp_path):
    df = pd.DataFrame(
        {
            "TrackingErr": [0.02],
            "Beta": [1.0],
            "AnnReturn": [0.05],
            "Capital": [100],
            "AnnVol": [0.02],
            "Agent": ["A"],
        }
    )
    assert isinstance(beta_te_scatter.make(df), go.Figure)

    arr = np.random.normal(size=(5, 4))
    assert isinstance(boxplot.make(arr), go.Figure)

    fig1 = risk_return.make(df.assign(ShortfallProb=[0.01]))
    fig2 = sharpe_ladder.make(df)
    cross = crossfilter.make([fig1, fig2], df)
    assert isinstance(cross, go.Figure)

    gp = grid_panel.make([fig1, fig2])
    hover_sync.apply([gp])
    assert isinstance(gp, go.Figure)

    tpl = dashboard_templates.get("default")
    assert isinstance(tpl, dict)

    grid_a = pd.DataFrame(
        {
            "AE_leverage": [1, 2],
            "ExtPA_frac": [0.2, 0.4],
            "Sharpe": [0.4, 0.5],
        }
    )
    grid_b = pd.DataFrame(
        {
            "AE_leverage": [1, 2],
            "ExtPA_frac": [0.2, 0.4],
            "Sharpe": [0.45, 0.55],
        }
    )
    assert isinstance(delta_heatmap.make(grid_a, grid_b), go.Figure)

    expo = pd.DataFrame({"Region": ["United States"], "Exposure": [1.0]})
    assert isinstance(geo_exposure.make(expo), go.Figure)

    assert isinstance(funnel.make(arr), go.Figure)

    assert isinstance(
        factor_timeline.make(pd.DataFrame({"F1": [0.1, 0.2]}, index=[0, 1])), go.Figure
    )

    base_fig = fan.make(arr)
    events = [(1, "A"), (2, "B")]
    assert isinstance(milestone_timeline.make(events, base_fig), go.Figure)

    assert isinstance(mosaic.make(arr), go.Figure)

    out = tmp_path / "report.pdf"
    pdf_report.save([base_fig], out)
    assert out.exists()

    assert isinstance(quantile_band.make(arr), go.Figure)

    table = rank_table.make(df)
    assert isinstance(table, pd.DataFrame)

    assert isinstance(scenario_play.make(arr), go.Figure)

    idx = pd.MultiIndex.from_product([[2020, 2021], [0, 1, 2, 3]], names=["Year", "Month"])
    df_season = pd.DataFrame({"R": arr[:2].reshape(-1)}, index=idx)
    assert isinstance(seasonality_heatmap.make(df_season), go.Figure)

    df_series = pd.DataFrame({"A": [1, 2, 3]})
    assert isinstance(spark_matrix.make(df_series), go.Figure)

    assert isinstance(surface_animation.make(grid_a), go.Figure)
    assert isinstance(surface_slice.make(grid_a), go.Figure)

    assert isinstance(triple_scatter.make(df), go.Figure)

    stack_df = pd.DataFrame({"A": [1, 2], "B": [2, 3]})
    assert isinstance(weighted_stack.make(stack_df), go.Figure)


def test_theme_reload(tmp_path):
    cfg = tmp_path / "theme.yaml"
    cfg.write_text(
        "colorway: ['#111111', '#222222']\n"
        "font: DejaVu Sans\n"
        "paper_bgcolor: '#eeeeee'\n"
        "plot_bgcolor: '#dddddd'\n"
    )
    from pa_core.viz import theme

    theme.reload_theme(cfg)
    assert theme.TEMPLATE.layout.font.family == "DejaVu Sans"  # type: ignore[attr-defined]
    assert list(theme.TEMPLATE.layout.colorway)[:2] == ["#111111", "#222222"]  # type: ignore[attr-defined]
    assert theme.TEMPLATE.layout.paper_bgcolor == "#eeeeee"  # type: ignore[attr-defined]
    assert theme.TEMPLATE.layout.plot_bgcolor == "#dddddd"  # type: ignore[attr-defined]
