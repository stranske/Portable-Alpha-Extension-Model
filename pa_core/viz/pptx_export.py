from __future__ import annotations

import base64
import io
import os
from pathlib import Path
from typing import Any, Iterable, Sequence, cast

import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches

from .export_backend import figure_to_png_bytes, is_browser_runtime, run_with_browser_png_cache

__all__ = ["render_chart_png", "add_chart_slide", "save", "save_async"]

# Tiny 1x1 PNG used as a placeholder in pytest or explicit placeholder mode so
# chart slides never spawn a potentially hanging kaleido/Chromium subprocess.
_ONE_PX_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMA"
    "ASsJTYQAAAAASUVORK5CYII="
)
_RENDERER_ERROR_TERMS = ("kaleido", "chrome", "chromium")


def _is_static_renderer_error(exc: BaseException) -> bool:
    return any(term in str(exc).lower() for term in _RENDERER_ERROR_TERMS)


def render_chart_png(fig: Any) -> bytes:
    """Render a Plotly figure to PNG bytes for embedding in a PPTX slide.

    In pytest or explicit placeholder mode a tiny placeholder image is returned
    to avoid hanging kaleido subprocesses. Otherwise a real static render is
    attempted via Kaleido; if that renderer is unavailable the failure is
    surfaced as an actionable ``RuntimeError`` rather than being silently
    swallowed (a silently-skipped chart produces a divergent, incomplete board
    pack).
    """
    if not is_browser_runtime() and (
        os.environ.get("PYTEST_CURRENT_TEST") or os.environ.get("PA_PPTX_PLACEHOLDER") == "1"
    ):
        return _ONE_PX_PNG
    try:
        return cast(bytes, figure_to_png_bytes(fig))
    except Exception as e:  # pragma: no cover - exercised only without kaleido
        if not _is_static_renderer_error(e):
            raise
        raise RuntimeError(
            "PPTX export requires a static image renderer (Kaleido/Chromium). "
            "Install Plotly Kaleido or Chrome/Chromium. For Debian/Ubuntu: "
            "pip install 'plotly[kaleido]' or sudo apt-get install -y chromium-browser. "
            f"Original error: {e}"
        ) from e


def add_chart_slide(
    prs: Any,
    fig: Any,
    *,
    alt: str | None = None,
    layout_index: int = 5,
) -> None:
    """Add a single full-bleed chart slide to ``prs`` for ``fig``.

    The image is rendered via :func:`render_chart_png` (shared error handling)
    and accessibility alt text is set from ``alt`` or the figure's layout title.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    img_bytes = render_chart_png(fig)
    pic = slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )
    if not alt:
        layout = getattr(fig, "layout", None)
        title = getattr(layout, "title", None) if layout is not None else None
        text = getattr(title, "text", None) if title is not None else None
        alt = str(text) if text else ""
    if alt:
        el = pic._element.xpath("./p:nvPicPr/p:cNvPr")[0]
        el.set("descr", alt)


def save(
    figs: Iterable[go.Figure],
    path: str | Path,
    *,
    alt_texts: Sequence[str] | None = None,
) -> None:
    """Save figures to a PowerPoint file, one per slide.

    Parameters
    ----------
    figs:
        Iterable of Plotly figures.
    path:
        Destination ``.pptx`` file.
    alt_texts:
        Optional sequence of alt text strings for accessibility. If omitted,
        each figure's layout title is used when present.

    Raises
    ------
    RuntimeError
        If a static image renderer (Kaleido/Chromium) is required but
        unavailable. This matches the committee export-packet behaviour so a
        missing renderer can never silently drop chart slides.
    """
    pres = Presentation()
    alt_iter = iter(alt_texts) if alt_texts is not None else None
    for fig in figs:
        add_chart_slide(pres, fig, alt=next(alt_iter, None) if alt_iter else None)
    pres.save(str(path))


async def save_async(
    figs: Iterable[go.Figure],
    path: str | Path,
    *,
    alt_texts: Sequence[str] | None = None,
) -> None:
    """Async browser-safe variant of :func:`save`.

    In server Python this delegates directly to the synchronous implementation.
    In stlite/Pyodide it first renders every figure through Plotly.js, activates
    the PNG cache, and then runs the same synchronous PPTX assembly.
    """
    figs_list = list(figs)
    await run_with_browser_png_cache(
        figs_list,
        lambda: save(figs_list, path, alt_texts=alt_texts),
    )
