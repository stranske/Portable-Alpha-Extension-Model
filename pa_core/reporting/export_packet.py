"""Committee-ready export packet (PPTX + Excel).

Used by the CLI and dashboard to generate a PowerPoint deck and a
companion Excel workbook. Provides alt-text support and actionable
errors when static chart export is unavailable.
"""

from __future__ import annotations

import base64
import io
import os
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence, Tuple

import pandas as pd
from pptx import Presentation as _Presentation  # type: ignore
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .excel import export_to_excel, finalize_excel_workbook

__all__ = ["create_export_packet"]

_ONE_PX_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMA"
    "ASsJTYQAAAAASUVORK5CYII="
)


def _add_title_slide(prs: Any, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_summary_table_slide(prs: Any, df: pd.DataFrame, title: str = "Executive Summary") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Title
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tx_box.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True

    head = df.copy()
    max_rows = 10
    if len(head) > max_rows:
        head = head.iloc[:max_rows]

    rows, cols = head.shape
    left, top = Inches(0.5), Inches(1.1)
    width, height = Inches(9), Inches(5)
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Header
    for c, name in enumerate(head.columns):
        cell = table.cell(0, c)
        cell.text = str(name)
        cell.text_frame.paragraphs[0].runs[0].font.bold = True

    # Body
    for r in range(rows):
        for c in range(cols):
            val = head.iat[r, c]
            cell = table.cell(r + 1, c)
            cell.text = f"{val:.2%}" if isinstance(val, float) and 0 <= val <= 1 else str(val)
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)


def _add_chart_slide(prs: Any, fig: Any, alt: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if os.environ.get("CI") or os.environ.get("PYTEST_CURRENT_TEST"):
        # Avoid hanging kaleido subprocesses in CI/pytest by using a tiny placeholder image.
        img = _ONE_PX_PNG
    else:
        try:
            img = fig.to_image(format="png", engine="kaleido")
        except Exception as e:
            raise RuntimeError(
                "PPTX export requires a static image renderer (Kaleido/Chromium). "
                "Install Plotly Kaleido or Chrome/Chromium. For Debian/Ubuntu: "
                "pip install 'plotly[kaleido]' or sudo apt-get install -y chromium-browser. "
                f"Original error: {e}"
            ) from e

    pic = slide.shapes.add_picture(io.BytesIO(img), Inches(0), Inches(0))
    if not alt:
        layout = getattr(fig, "layout", None)
        title = getattr(layout, "title", None) if layout is not None else None
        text = getattr(title, "text", None) if title is not None else None
        alt = str(text) if text else ""
    if alt:
        el = pic._element.xpath("./p:nvPicPr/p:cNvPr")[0]
        el.set("descr", alt)


def _add_table_slide(prs: Any, df: pd.DataFrame, title: str = "Table") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Title
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tx_box.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True

    head = df.copy()
    max_rows = 12
    if len(head) > max_rows:
        head = head.iloc[:max_rows]

    rows, cols = head.shape
    left, top = Inches(0.5), Inches(1.1)
    width, height = Inches(9), Inches(5)
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Header
    for c, name in enumerate(head.columns):
        cell = table.cell(0, c)
        cell.text = str(name)
        cell.text_frame.paragraphs[0].runs[0].font.bold = True

    # Body
    for r in range(rows):
        for c in range(cols):
            val = head.iat[r, c]
            cell = table.cell(r + 1, c)
            cell.text = f"{val:.3f}" if isinstance(val, float) else str(val)
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)


def create_export_packet(
    *,
    figs: Iterable[Any],
    summary_df: pd.DataFrame,
    raw_returns_dict: dict[str, pd.DataFrame],
    inputs_dict: dict[str, object],
    base_filename: str | Path = "committee_packet",
    alt_texts: Sequence[str] | None = None,
    pivot: bool = False,
    manifest: Mapping[str, Any] | None = None,
    prev_summary_df: pd.DataFrame | None = None,
    prev_manifest: Mapping[str, Any] | None = None,
    stress_delta_df: pd.DataFrame | None = None,
) -> Tuple[str, str]:
    """Create PPTX + Excel packet and return their paths."""
    base = Path(str(base_filename))
    pptx_path = str(base.with_suffix(".pptx"))
    excel_path = str(base.with_suffix(".xlsx"))

    from .run_diff import build_run_diff

    cfg_diff_df: pd.DataFrame | None = None
    metric_diff_df: pd.DataFrame | None = None
    if prev_manifest is not None or prev_summary_df is not None:
        prev_sum = prev_summary_df if prev_summary_df is not None else pd.DataFrame()
        cfg_diff_df, metric_diff_df = build_run_diff(manifest, prev_manifest, summary_df, prev_sum)

    # Excel workbook (full tables)
    finalize_after_append = stress_delta_df is not None and not stress_delta_df.empty
    export_to_excel(
        inputs_dict,
        summary_df,
        raw_returns_dict,
        filename=excel_path,
        pivot=pivot,
        diff_config_df=cfg_diff_df,
        diff_metrics_df=metric_diff_df,
        finalize=not finalize_after_append,
    )
    if finalize_after_append:
        with pd.ExcelWriter(
            excel_path, engine="openpyxl", mode="a", if_sheet_exists="replace"
        ) as writer:
            stress_delta_df.to_excel(writer, sheet_name="StressDelta", index=False)
        finalize_excel_workbook(excel_path, inputs_dict, summary_df)

    # PowerPoint deck
    prs = _Presentation()
    _add_title_slide(
        prs,
        title="Portfolio Analysis Packet",
        subtitle=f"Generated: {pd.Timestamp.now():%Y-%m-%d %H:%M}",
    )
    if not summary_df.empty:
        _add_summary_table_slide(prs, summary_df)

    figs_list = list(figs)
    try:
        sens_val = inputs_dict.get("_sensitivity_df")
        sens_df: pd.DataFrame | None = sens_val if isinstance(sens_val, pd.DataFrame) else None
        if sens_df is not None and not sens_df.empty:
            if {"Parameter", "DeltaAbs"} <= set(sens_df.columns):
                from ..viz import tornado

                series = tornado.series_from_sensitivity(sens_df)
                title = "Sensitivity Tornado"
                has_tornado = False
                for fig in figs_list:
                    layout = getattr(fig, "layout", None)
                    fig_title = getattr(layout, "title", None) if layout else None
                    text = getattr(fig_title, "text", None) if fig_title else None
                    if text and str(text).strip().lower() == title.lower():
                        has_tornado = True
                        break
                if not has_tornado:
                    figs_list.append(tornado.make(series, title=title))
    except Exception:
        # Best-effort only; tornado chart is optional.
        pass

    alt_iter = iter(alt_texts) if alt_texts is not None else None
    for fig in figs_list:
        _add_chart_slide(prs, fig, next(alt_iter, None) if alt_iter else None)

    # Appendix reminding where detailed tables live
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = tx_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Appendix: Full tables are included in the Excel workbook."
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.LEFT

    # Diff appendix comparing to previous run
    if cfg_diff_df is not None and not cfg_diff_df.empty:
        _add_table_slide(prs, cfg_diff_df, title="Config Changes")
    if metric_diff_df is not None and not metric_diff_df.empty:
        _add_table_slide(prs, metric_diff_df, title="Metric Changes")
    if stress_delta_df is not None and not stress_delta_df.empty:
        from .stress_delta import format_delta_table_text

        _add_table_slide(
            prs,
            format_delta_table_text(stress_delta_df),
            title="Stress Delta vs Base",
        )

    # Optional manifest summary appendix
    if manifest:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
        tf = tx_box.text_frame
        tf.clear()
        # Title
        title_run = tf.paragraphs[0].add_run()
        title_run.text = "Reproducibility Manifest"
        title_run.font.size = Pt(20)
        title_run.font.bold = True

        def _add_line(text: str) -> None:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = text
            r.font.size = Pt(11)

        git_commit = str(manifest.get("git_commit", "unknown"))
        timestamp = str(manifest.get("timestamp", ""))
        seed = manifest.get("seed")
        data_files = manifest.get("data_files") or {}
        cli_args = manifest.get("cli_args") or {}
        cfg = manifest.get("config") or {}

        _add_line(f"Commit: {git_commit}")
        _add_line(f"Timestamp (UTC): {timestamp}")
        if seed is not None:
            _add_line(f"Seed: {seed}")
        if isinstance(data_files, dict) and data_files:
            _add_line("Data files (sha256):")
            for path, h in list(data_files.items())[:8]:
                _add_line(f"  • {Path(path).name}: {h[:12]}…")
            if len(data_files) > 8:
                _add_line(f"  • … and {len(data_files) - 8} more")
        if isinstance(cli_args, dict) and cli_args:
            mode = cli_args.get("mode")
            _add_line(f"Mode: {mode}")
        if isinstance(cfg, dict) and cfg:
            _add_line("Key config fields:")
            keys = [
                "N_SIMULATIONS",
                "N_MONTHS",
                "w_beta_H",
                "w_alpha_H",
            ]
            for k in keys:
                if k in cfg:
                    _add_line(f"  • {k}: {cfg[k]}")

    prs.save(pptx_path)
    return pptx_path, excel_path
